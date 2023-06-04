import shutil
import time
from datetime import datetime

from pptx import Presentation
import openai
import json
import os
import re
import asyncio

openai.api_key = os.environ.get('API_KEY')
CONTENT = [
    {"role": "system", "content": "Can you explain the slides in basic english, and provide examples if needed!"}
]
ERROR_MESSAGE = "Something is wrong:"
ENGINE_MODEL = "gpt-3.5-turbo"
WRITE_TO_FILE_MODE = 'w'
UPLOADS_FOLDER = 'uploads'
OUTPUTS_FOLDER = 'outputs'
PROCESSED_FOLDER = 'processed'


async def parse_presentation(presentation_path):
    """
    This method receives a path for a pptx presentation, checks if the path is found in the operating system, if so,
    parses the data to slides. The method, returns a list of explanations.
    :param: presentation_path: path of a power-point presentation. (String)
    :return: list of explanations. (List of strings)
    """
    # check if path is available
    if not os.path.isfile(presentation_path):
        print(f"{ERROR_MESSAGE} the path you provided does not exist.")
        return []

    prs = Presentation(presentation_path)
    explanations = []
    for slide in prs.slides:
        explanation = await parse_slide_of_pptx(slide)
        explanations.append(explanation)
    return explanations


async def parse_slide_of_pptx(slide):
    """
    This method receives a single slide, the method parses the data on that slide, it calls another method to get the
    response and return it to the parse_presentation method. It throws an error if an exception occurred.
    :param: slide: A single slide from the power-point. (slide object)
    :return: The explanation if the processing went well, an error, otherwise. (List of Strings)
    """
    try:
        response = await request_completion(slide)
        return response
    except Exception as error:
        error_message = f"{ERROR_MESSAGE} Error processing slide: {str(error)}"
        return error_message


def parse_text_of_slide(slide):
    """
    This method receives a slide, and extracts all the extractable text found in that slide. In addition, it cleans
    the data using the strip function.
    :param: slide: A single slide from the power-point. (Slide Object)
    :return: Explanation response. (String)
    """
    slide_text = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                slide_text.append(run.text.strip())
    return " ".join(slide_text)


async def request_completion(slide):
    """
    This method receives a slide object, sends a request to the openai API asking the server to explain the content
    of that slide, the method returns the response from the API.
    :param: slide: Slide: a single slide from the power-point. (Slide Object)
    :return: Response of the API.
    """
    CONTENT.append({"role": "user", "content": parse_text_of_slide(slide)})
    response = await openai.ChatCompletion.acreate(
        model=ENGINE_MODEL,
        messages=CONTENT
    )
    content = response["choices"][0].message.content
    cleaned_content = clean_text(content)
    return cleaned_content


def clean_text(text):
    """
    This method receives the response retrieved from the openai API and cleans it, in other words, gets rid of
    unwanted characters.
    :param: text: Response from API. (string)
    :return: clean Version of the response. (string)
    """
    cleaned_text = re.sub(r"\n", "", text)
    cleaned_text = cleaned_text.encode("ascii", "ignore").decode("utf-8")
    return cleaned_text.strip()


def save_explanations(explanations, file_path):
    """
    This method receives a list of strings, each string representing an explanation. It creates a JSON file and appends
    the explanations to the file, using the original file name.
    :param: explanations: List of explanations retrieved from the API. (List of strings)
    :param: file_path: Path of the original file. (String)
    :return:
    """
    file_name = os.path.basename(file_path)
    presentation_name, extension = os.path.splitext(file_name)
    uid = re.search(r"\w{8}-\w{4}-\w{4}-\w{4}-\w{12}", presentation_name).group(0)
    print(presentation_name)
    output_file = os.path.join(OUTPUTS_FOLDER, f"{presentation_name}_explanations.json")
    slide_explanations = {}

    for slide_num, explanation in enumerate(explanations, start=1):
        slide_key = f"slide{slide_num}"
        slide_explanations[slide_key] = explanation

    try:
        with open(output_file, "w") as file:
            json.dump(slide_explanations, file, indent=4)
        print(f"Explanations saved to {output_file}")
    except IOError as error:
        print(f"Error saving explanations: {str(error)}")


def move_file(file_path, destination_folder):
    """
    This method receives a file path, and a destination folder. The method moves the processed folder from the
    'uploads' folder to the 'processed' folder in order to prevent the explainer to scan old files that have been
    already processed
    :param file_path: path of a file (string)
    :param destination_folder: name of the processed folder (string)
    :return:
    """
    file_name = os.path.basename(file_path)
    destination_path = os.path.join(destination_folder, file_name)
    shutil.move(file_path, destination_path)
    print(f"Moved file: {file_path} to {destination_path}")


async def process_file(file_path):
    """
    This method receives a file_path, calls the needed functions to explain the contents of the power-point, save and
    moves the processed file. it throws an exception if it could not process a file.
    :param file_path: a file path from the 'uploads' folder (string)
    :return:
    """
    print(f"Processing file: {file_path}")
    try:
        explanations = await parse_presentation(file_path)
        save_explanations(explanations, file_path)
        move_file(file_path, PROCESSED_FOLDER)
    except Exception as error:
        error_message = f"{ERROR_MESSAGE} Error processing file: {str(error)}"
        print(error_message)


async def main_loop():
    """
    This method keeps running in an infinite loop, each iteration it scans the 'uploads' folder for a new file to
    process, explains the content of the file, and then sleeps for 10 seconds
    :return:
    """
    while True:
        for file_name in os.listdir(UPLOADS_FOLDER):
            file_path = os.path.join(UPLOADS_FOLDER, file_name)
            if os.path.isfile(file_path):
                await process_file(file_path)

        await asyncio.sleep(10)


if __name__ == "__main__":
    print("Explainer started.")
    asyncio.run(main_loop())
